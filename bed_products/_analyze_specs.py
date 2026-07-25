"""分析PPT规格组件名 → 建议映射到价格表规格名"""
import os, re
from pptx import Presentation

BASE_DIR = r'G:\Trae工作文件\顾家产品库'
SOFA_DIR = os.path.join(BASE_DIR, 'sofa_products')

def get_pptx_texts(p):
    t = []
    try:
        for s in Presentation(p).slides:
            for sh in s.shapes:
                if sh.has_text_frame:
                    for pa in sh.text_frame.paragraphs:
                        x = pa.text.strip()
                        if x: t.append(x)
    except:
        pass
    return t

# 各产品在价格表中的规格名（从之前分析得来）
price_spec_examples = {
    '9659': ['1.5右','1.5左','1.5无组','1双','2.5右','2.5左','3右','3左','1右','1左'],
    'JD.0061': ['1.5右','1.5右A','1.5左','1.5左A','1右','1左','3右','3左','3.5右','3.5左'],
    'JD.0036': ['1.5右','1.5左','1.5无组','1双','3右','3左'],
    'JD.0069': ['1.5右','1.5左','1.5无组','1双'],
    'JD.0072': ['1.5右','1.5左','1.5无组','1双','3右','3左'],
    'JD.0077': ['1.5右','1.5左','1.5无组','1双'],
    'JD.0006': ['1.5右','1.5左','1.5无组','3右','3左','1右','1左'],
    'JD.0020': ['1.5右','1.5左','1.5无组','3右','3左'],
    'JD.6013': ['1.5右','1.5左','1.5无组','1.5右电动','1.5左电动'],
    'JD.6016': ['1.5右','1.5左','1.5无组','1.5右电动','1.5左电动'],
    'JD.6025': ['1.5右','1.5左','1.5无组','1.5右电动','1.5左电动'],
    'JD.0062': ['1.5右','1.5左','1.5无组','3右','3左'],
    'JD.0006B': ['1.5右','1.5左','1.5无组','3右','3左','1右','1左'],
}

def clean_ppt_comp(name):
    """Clean a PPT component name for matching"""
    n = name.strip()
    n = n.replace('（', '(').replace('）', ')')
    # Remove suffixes like (坐宽85)
    n = re.sub(r'[（(].+?[）)]', '', n).strip()
    return n

print("=== 各产品PPT组件 → 建议映射到价格表规格名 ===\n")

for folder in sorted(os.listdir(SOFA_DIR)):
    fp = os.path.join(SOFA_DIR, folder)
    if not os.path.isdir(fp) or folder.startswith('~$'):
        continue
    ppt_dir = os.path.join(fp, 'PPT')
    if not os.path.isdir(ppt_dir):
        continue
    fs = os.listdir(ppt_dir)
    pptx = [f for f in fs if f.endswith('.pptx')]
    if not pptx:
        continue
    
    texts = get_pptx_texts(os.path.join(ppt_dir, pptx[0]))
    comps_set = set()
    for t in texts:
        tc = t.replace('\n', ' ').strip()
        if not re.search(r'\d+(?:\.\d+)?\s*(?:CM|cm|m)\b', tc):
            continue
        if any(kw in tc for kw in ['由于手工测量', '误差', 'mm']):
            continue
        tc = tc.replace('丨', '|')
        desc = re.sub(r'^[A-Z]{1,4}\s*[\-—]\s*', '', tc).strip('-— ')
        m = re.search(r'(.+?)\s*[：:]\s*\d+', desc)
        if not m:
            m = re.search(r'(.+?)\s*[-—]\s*\d+\.?\d*\s*m', desc)
        if not m:
            m = re.search(r'(.+?)\s*\d{3,}\s*(?:CM|cm)', desc)
        if m:
            name = m.group(1).strip()
            parts = re.split(r'[\+\+｜|]', name.replace('，', '+'))
            for p in parts:
                p = p.strip()
                if p and len(p) > 1 and len(p) < 20:
                    comps_set.add(clean_ppt_comp(p))

    if not comps_set or folder not in price_spec_examples:
        continue

    pspecs = price_spec_examples[folder]
    print(f"\n{folder}:")
    for pc in sorted(comps_set):
        # Try fuzzy match
        best = None
        pc_clean = re.sub(r'[\s\-单/]', '', pc)
        for ps in pspecs:
            ps_clean = re.sub(r'[\s\-]', '', ps)
            if pc_clean == ps_clean:
                best = ps
                break
        if best:
            print(f"  ✓ '{pc}' → '{best}'")
        else:
            # Try removing trailing A
            pc2 = re.sub(r'A$', '', pc_clean)
            for ps in pspecs:
                ps2 = re.sub(r'A$', '', re.sub(r'[\s\-]', '', ps))
                if pc2 == ps2:
                    print(f"  ~ '{pc}' → '{ps}'  (去掉'A')")
                    best = ps
                    break
        if not best:
            # Try just the number part
            num = re.search(r'(\d+\.?\d*)', pc)
            if num:
                n = num.group(1)
                for ps in pspecs:
                    if n in ps and ('左' in pc and '左' in ps) or ('右' in pc and '右' in ps) or ('无' in pc and '无' in ps):
                        print(f"  ? '{pc}' → 可能是 '{ps}'")
                        best = ps
                        break
        if not best:
            print(f"  ✗ '{pc}' → 未匹配 (价格表中有: {', '.join(pspecs[:8])})")
