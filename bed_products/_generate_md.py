import os, json, re
from pptx import Presentation
from pypdf import PdfReader
import openpyxl

BASE_DIR = r'G:\Trae工作文件\顾家产品库'
OUTPUT_DIR = os.path.join(BASE_DIR, 'markdown_db')
BED_DIR = os.path.join(BASE_DIR, 'bed_products')
MATTRESS_DIR = os.path.join(BASE_DIR, 'mattress_products')
SOFA_DIR = os.path.join(BASE_DIR, 'sofa_products')
JSON_PATH = os.path.join(OUTPUT_DIR, '_price_data.json')

os.makedirs(OUTPUT_DIR, exist_ok=True)

# Load price data
price_data = {'bed_frames': [], 'mattresses': []}
if os.path.exists(JSON_PATH):
    with open(JSON_PATH, 'r', encoding='utf-8') as f:
        price_data = json.load(f)

frame_by_code = {}
for item in price_data['bed_frames']:
    code = item['货号']
    frame_by_code.setdefault(code, []).append(item)

mattress_by_code = {}
for item in price_data['mattresses']:
    code = item['货号']
    mattress_by_code.setdefault(code, []).append(item)

def normalize(s):
    return re.sub(r'[\s.\n]', '', s.upper())

def match_code(folder, code_dict):
    fn = normalize(folder)
    for code, items in code_dict.items():
        if normalize(code) == fn:
            return items
    fs = re.sub(r'^(JD|HS|BY)', '', fn)
    for code, items in code_dict.items():
        cs = re.sub(r'^(JD|B|90|HS|BY)', '', normalize(code))
        if fs and fs == cs:
            return items
        for n in re.findall(r'\d+[A-Z0-9]*', fs):
            if len(n) >= 4 and n in cs:
                return items
    return []

# ==================== PPT/PDF Text Extraction ====================
def get_pptx_texts(pptx_path):
    texts = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
    except Exception as e:
        print(f'  [WARN] pptx: {e}')
    return texts

def get_pdf_texts(pdf_path):
    texts = []
    try:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                texts.extend([l.strip() for l in t.split('\n') if l.strip()])
    except Exception as e:
        print(f'  [WARN] pdf: {e}')
    return texts

# ==================== Sofa: Extract Style & Dimensions ====================
def extract_sofa_style_dimensions(texts):
    """从PPT/PDF文本中提取沙发适配风格和尺寸"""
    result = {'风格': '', '尺寸表': {}}
    all_text = '\n'.join(texts)
    
    # --- 适配风格 ---
    style_keywords = ['风格', '适配', '现代', '意式', '奶油', '极简', '简约', '轻奢',
                      '复古', '中古', '北欧', '法式', '轻中古']
    style_sentences = []
    for t in texts:
        t_clean = t.replace('\n', '')
        if any(kw in t_clean for kw in style_keywords):
            if len(t_clean) < 200:  # short enough
                style_sentences.append(t_clean)
    
    # Deduplicate and build style desc
    unique_styles = []
    for s in style_sentences:
        # Skip common non-style text
        if any(kw in s for kw in ['尺寸', '规格', 'CM', 'cm', 'mm', '介绍', '概览']):
            continue
        if s not in unique_styles:
            unique_styles.append(s)
    
    if unique_styles:
        result['风格'] = '；'.join(unique_styles[:5])
    
    # --- 提取尺寸 ---
    # 多种PPT尺寸格式：
    # 格式1: "S-3单A+1.5单A：292CM"
    # 格式2: "S: 2.5左+1右  2.72m"
    # 格式3: "M-3双 287cm"
    raw_dims = []
    dim_table = {}
    
    for t in texts:
        t_clean = t.replace('\n', ' ').strip()
        
        # 格式1: 规格名-描述：数字CM (e.g. "S-3单A+1.5单A：292CM")
        m1 = re.findall(r'([A-Z]+[\-\u4e00-\u9fff\w/\+]+)\s*[：:]\s*(\d+(?:\.\d+)?)\s*(?:CM|cm|Cm)', t_clean)
        for name, size in m1:
            # Extract the spec letter
            spec = re.match(r'([A-Z]+)', name)
            if spec:
                dim_table[spec.group(1)] = f"{size}cm"
        
        # 格式2: 规格名: 描述 数字m (e.g. "S: 2.5左+1右  2.72m")
        m2 = re.findall(r'([A-Z]+)\s*[：:]\s*(?:[\u4e00-\u9fff\w\+/]+)\s*(\d+(?:\.\d+)?)\s*[mM]', t_clean)
        for spec, size_m in m2:
            size_cm = str(int(float(size_m) * 100)) if float(size_m) < 10 else size_m
            dim_table[spec.upper()] = f"{size_cm}cm"
        
        # 格式3: 数字CM直接出现 (e.g. "L-大3双 311cm")
        m3 = re.findall(r'([A-Z]+)\s*[\-\u4e00-\u9fff\w/\+]*\s*(\d+(?:\.\d+)?)\s*(?:CM|cm|Cm)', t_clean)
        for spec, size in m3:
            if spec.upper() not in dim_table:
                dim_table[spec.upper()] = f"{size}cm"
        
        # 收集所有尺寸行
        if re.search(r'\d{2,4}\s*[xX×*]\s*\d{2,4}', t_clean) or re.search(r'\d+(?:\.\d+)?\s*(?:[cC][mM]|[mM])\b', t_clean):
            raw_dims.append(t_clean)
    
    result['尺寸表'] = dim_table
    result['raw_dimensions'] = raw_dims[:10]
    
    return result

# ==================== Bed: Extract 床架尺寸 (L*W*H) ====================
def extract_bed_dimensions(xlsx_path):
    """从床架话术xlsx中提取床架尺寸(L*W*H CM)和在售规格"""
    dims = {'床架尺寸': '', '在售规格': ''}
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        ws = wb.active
        for r in range(1, ws.max_row + 1):
            c2 = str(ws.cell(r, 2).value or '').strip()
            c3 = str(ws.cell(r, 3).value or '').strip()
            
            if '床架尺寸' in c2 and c3:
                dims['床架尺寸'] = c3.replace('\\', '\n').strip()
            elif '在售规格' in c2 and c3:
                dims['在售规格'] = c3.replace('\\', '\n').strip()
        wb.close()
    except Exception as e:
        print(f'    [WARN] 床架尺寸读取: {e}')
    return dims

# ==================== Utility ====================
def safe_name(folder):
    return folder.replace(' ', '_').replace('/', '_')

def write_md(md_path, lines):
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))

def price_str(val):
    if val == 0: return ''
    return f"¥{val:,}" if isinstance(val, int) else f"¥{val:,.0f}"

def list_images(img_dir):
    if not os.path.isdir(img_dir): return []
    exts = ('.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp')
    return sorted([f for f in os.listdir(img_dir) if f.lower().endswith(exts)])

def get_images(folder_path):
    result = {}
    for sub in ['场景图', '浏览图', '入户实景图', '白底图']:
        imgs = list_images(os.path.join(folder_path, sub))
        if imgs:
            result[sub] = imgs
    return result

def format_images(images_dict, folder_path):
    lines = []
    for subdir, imgs in sorted(images_dict.items()):
        lines.append(f"\n### {subdir} ({len(imgs)}张)")
        for img in imgs[:10]:
            rel = os.path.relpath(os.path.join(folder_path, subdir, img), BASE_DIR).replace('\\', '/')
            lines.append(f"![]({rel})")
        if len(imgs) > 10:
            lines.append(f"... 共{len(imgs)}张")
    return lines

def read_doc(folder_path):
    """读取产品文件夹中的文档"""
    ppt_dir = os.path.join(folder_path, 'PPT')
    if not os.path.isdir(ppt_dir): return None
    files = os.listdir(ppt_dir)
    xlsx = [f for f in files if f.endswith('.xlsx') and not f.startswith('~$')]
    pptx = [f for f in files if f.endswith('.pptx')]
    pdf = [f for f in files if f.endswith('.pdf')]
    png = [f for f in files if f.lower().endswith('.png') and '详情' not in f]
    jpg = [f for f in files if f.lower().endswith(('.jpg', '.jpeg')) and '详情' not in f]
    
    if xlsx:
        return {'type': 'xlsx', 'file': xlsx[0], 'path': os.path.join(ppt_dir, xlsx[0])}
    elif pptx:
        return {'type': 'pptx', 'file': pptx[0], 'path': os.path.join(ppt_dir, pptx[0])}
    elif pdf:
        return {'type': 'pdf', 'file': pdf[0], 'path': os.path.join(ppt_dir, pdf[0])}
    elif png:
        return {'type': 'image', 'files': png, 'path': ppt_dir}
    elif jpg:
        return {'type': 'image', 'files': jpg, 'path': ppt_dir}
    return None

# ==================== Generate ====================
print("=" * 60)

# ---------- 1. bed_products ----------
print("\n>>> bed_products")
for folder in sorted(os.listdir(BED_DIR)):
    fp = os.path.join(BED_DIR, folder)
    if not os.path.isdir(fp) or folder.startswith('~$'): continue
    
    doc = read_doc(fp)
    images = get_images(fp)
    frames = match_code(folder, frame_by_code)
    mattresses = match_code(folder, mattress_by_code)
    
    lines = [f"# {folder}\n"]
    
    # --- 产品信息 ---
    if doc:
        lines.append("## 产品信息\n")
        if doc['type'] == 'xlsx':
            info = {}
            bed_dims = extract_bed_dimensions(doc['path'])
            try:
                wb = openpyxl.load_workbook(doc['path'], data_only=True)
                ws = wb.active
                for r in range(1, ws.max_row + 1):
                    c1 = str(ws.cell(r,1).value or '').replace('\n','').strip()
                    c2 = str(ws.cell(r,2).value or '').strip()
                    c3 = str(ws.cell(r,3).value or '').strip()
                    if not c2 or c2.startswith('='): continue
                    if '设计风格' in c2: info['设计风格'] = c3[:200]
                    elif '产品配色' in c2 or '配色' in c2: info['产品配色'] = c3[:200]
                    elif ('材质' in c2 or '面料' in c2 or '靠包' in c2) and '新材质' not in c1:
                        if c3: info['材质'] = c3[:300].split('\n')[0].strip()
                    elif ('一句话' in c2 or '卖点' in c1 or '核心' in c1) and '新材质' not in c1:
                        if c3 and len(c3) > 5: info['卖点'] = c3[:500]
                wb.close()
            except: pass
            
            for k, v in info.items():
                if k != '卖点' and v:
                    lines.append(f"- **{k}**: {v}")
            if info.get('卖点'):
                lines.append(f"\n### 核心卖点\n{info['卖点']}")
            
            # 床架详细尺寸
            if bed_dims['床架尺寸']:
                lines.append(f"\n### 床架尺寸 (CM)")
                lines.append("> 尺寸格式: 床头到床尾距离 × 床宽 × 床头高度")
                for size in bed_dims['床架尺寸'].split('\n'):
                    size = size.strip()
                    if size:
                        # Try to parse L*W*H
                        parts = re.split(r'[*×xX]', size)
                        if len(parts) == 3:
                            lines.append(f"- **{size}** → 长度{parts[0].strip()}cm, 宽度{parts[1].strip()}cm, 床头高度{parts[2].strip()}cm")
                        else:
                            lines.append(f"- {size}")
            if bed_dims['在售规格']:
                lines.append(f"\n**在售规格**: {bed_dims['在售规格']}")
        lines.append("")
    
    # --- 床架价格 ---
    if frames:
        lines.append("## 床架\n")
        lines.append("> 规格格式: 长度(床头到床尾) × 床宽。例如 **201×151** = 长度201cm × 床宽151cm")
        by_color = {}
        for item in frames:
            by_color.setdefault(item.get('配色型号',''), []).append(item)
        for color, items in by_color.items():
            if color: lines.append(f"### 配色型号: {color}\n")
            lines.append("| 规格(CM) | 实际成交价 | 可搭配床垫厚度 |")
            lines.append("|----------|-----------|--------------|")
            for item in items:
                lines.append(f"| {item['规格']} | {price_str(item['实际成交价'])} | {item.get('可搭配床垫厚度','')} |")
            lines.append("")
    
    # --- 床垫价格 ---
    if mattresses:
        lines.append("## 床垫\n")
        s = mattresses[0].get('产品系列','')
        if s: lines.append(f"**产品系列**: {s}\n")
        m = mattresses[0].get('材质','')
        if m: lines.append(f"**材质**: {m}\n")
        lines.append("| 货号 | 规格 | 实际成交价 |")
        lines.append("|------|------|-----------|")
        for item in mattresses:
            lines.append(f"| {item['货号']} | {item['规格']} | {price_str(item['实际成交价'])} |")
        lines.append("")
    
    # --- 图片 ---
    if images:
        lines.append("## 图片素材\n")
        lines.extend(format_images(images, fp))
        lines.append("")
    
    write_md(os.path.join(OUTPUT_DIR, f"{safe_name(folder)}.md"), lines)
    ic = sum(len(v) for v in images.values())
    print(f"  {folder}: {len(frames)}床架 {len(mattresses)}床垫 {ic}图")

# ---------- 2. mattress_products ----------
print("\n>>> mattress_products")
for folder in sorted(os.listdir(MATTRESS_DIR)):
    fp = os.path.join(MATTRESS_DIR, folder)
    if not os.path.isdir(fp) or folder.startswith('~$'): continue
    
    doc = read_doc(fp)
    images = get_images(fp)
    matched = match_code(folder, mattress_by_code)
    
    lines = [f"# {folder}\n"]
    
    if doc and doc['type'] == 'xlsx':
        lines.append("## 产品信息\n")
        info = {}
        try:
            wb = openpyxl.load_workbook(doc['path'], data_only=True)
            ws = wb.active
            for r in range(1, ws.max_row + 1):
                c1 = str(ws.cell(r,1).value or '').replace('\n','').strip()
                c2 = str(ws.cell(r,2).value or '').strip()
                c3 = str(ws.cell(r,3).value or '').strip()
                if not c2 or c2.startswith('='): continue
                if ('材质' in c2 or '面料' in c2) and '新材质' not in c1:
                    if c3: info['材质'] = c3[:300].split('\n')[0].strip()
                elif '规格' in c2 or '尺寸' in c2:
                    info['规格'] = c3[:300]
                elif '产品配置' in c2:
                    info['产品配置'] = c3[:500]
                elif '睡感' in c2 and '等级' not in c2:
                    info['睡感'] = c3[:300]
                elif ('一句话' in c2 or '卖点' in c1 or '核心' in c1) and '新材质' not in c1:
                    if c3 and len(c3) > 5: info['卖点'] = c3[:500]
            wb.close()
        except: pass
        for k, v in info.items():
            if k != '卖点' and v: lines.append(f"- **{k}**: {v}")
        if info.get('卖点'):
            lines.append(f"\n### 核心卖点\n{info['卖点']}")
        lines.append("")
    
    if matched:
        lines.append("## 床垫\n")
        s = matched[0].get('产品系列','')
        if s: lines.append(f"**产品系列**: {s}\n")
        m = matched[0].get('材质','')
        if m: lines.append(f"**材质**: {m}\n")
        lines.append("| 货号 | 规格 | 实际成交价 |")
        lines.append("|------|------|-----------|")
        for item in matched:
            lines.append(f"| {item['货号']} | {item['规格']} | {price_str(item['实际成交价'])} |")
        lines.append("")
    
    if images:
        lines.append("## 图片素材\n")
        lines.extend(format_images(images, fp))
        lines.append("")
    
    write_md(os.path.join(OUTPUT_DIR, f"{safe_name(folder)}.md"), lines)
    ic = sum(len(v) for v in images.values())
    print(f"  {folder}: {len(matched)}床垫 {ic}图")

# ---------- 3. sofa_products ----------
print("\n>>> sofa_products")
for folder in sorted(os.listdir(SOFA_DIR)):
    fp = os.path.join(SOFA_DIR, folder)
    if not os.path.isdir(fp) or folder.startswith('~$'): continue
    
    doc = read_doc(fp)
    images = get_images(fp)
    
    lines = [f"# {folder}\n"]
    
    # --- 产品信息 ---
    style_info = {}
    if doc and doc['type'] in ('pptx', 'pdf'):
        texts = get_pptx_texts(doc['path']) if doc['type'] == 'pptx' else get_pdf_texts(doc['path'])
        
        lines.append("## 产品信息\n")
        lines.append(f"**PPT文档**: {doc['file']}\n")
        
        if not texts:
            lines.append("> 该文档为图片型文件，请直接查看源文件获取产品详情。\n")
        else:
            style_info = extract_sofa_style_dimensions(texts)
            
            # 适配风格
            if style_info.get('风格'):
                lines.append(f"### 适配风格\n{style_info['风格']}\n")
            
            # 规格尺寸表
            dim_table = style_info.get('尺寸表', {})
            raw_dims = style_info.get('raw_dimensions', [])
            
            if dim_table:
                lines.append("### 规格尺寸\n")
                lines.append("| 规格 | 尺寸(CM) |")
                lines.append("|------|---------|")
                for name, size in dim_table.items():
                    lines.append(f"| {name} | {size} |")
                lines.append("")
            elif raw_dims:
                lines.append("### 规格尺寸\n")
                for d in raw_dims[:10]:
                    lines.append(f"- {d}")
                lines.append("")
            
            # Full text summary
            lines.append("### 内容概要\n")
            seen = set()
            count = 0
            for t in texts:
                t_clean = t.strip()
                if t_clean and t_clean not in seen and len(t_clean) > 5:
                    seen.add(t_clean)
                    if count < 20:
                        lines.append(f"- {t_clean[:120]}")
                    count += 1
            if count > 20:
                lines.append(f"- ... 共{count}条文本")
            lines.append("")
    elif doc and doc['type'] == 'xlsx':
        lines.append("## 产品信息\n")
        info = {}
        try:
            wb = openpyxl.load_workbook(doc['path'], data_only=True)
            ws = wb.active
            for r in range(1, ws.max_row + 1):
                c2 = str(ws.cell(r,2).value or '').strip()
                c3 = str(ws.cell(r,3).value or '').strip()
                if not c2 or c2.startswith('='): continue
                if '设计风格' in c2: info['设计风格'] = c3[:200]
                elif '产品配色' in c2 or '配色' in c2: info['产品配色'] = c3[:200]
            wb.close()
        except: pass
        for k, v in info.items():
            if v: lines.append(f"- **{k}**: {v}")
        lines.append("")
    elif doc and doc['type'] == 'image':
        lines.append("## 产品信息\n")
        lines.append("**参考文档**:\n")
        for img_file in doc['files']:
            rel = os.path.relpath(os.path.join(doc['path'], img_file), BASE_DIR).replace('\\', '/')
            lines.append(f"![]({rel})")
        lines.append("")
    
    # --- 图片 ---
    if images:
        lines.append("## 图片素材\n")
        lines.extend(format_images(images, fp))
        lines.append("")
    
    if not images and not doc:
        lines.append("*(暂无资料)*\n")
    
    write_md(os.path.join(OUTPUT_DIR, f"{safe_name(folder)}.md"), lines)
    ic = sum(len(v) for v in images.values())
    print(f"  {folder}: {ic}图")

print("\n=== 全部完成 ===")
