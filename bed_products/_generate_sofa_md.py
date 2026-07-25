"""沙发品类 - 从价格表取组合规格，从PPT找尺寸"""
import os, re, json
import win32com.client as win32
import pythoncom
from pptx import Presentation
from pypdf import PdfReader
import openpyxl

BASE_DIR = r'G:\Trae工作文件\顾家产品库'
OUTPUT_DIR = os.path.join(BASE_DIR, 'markdown_db')
SOFA_DIR = os.path.join(BASE_DIR, 'sofa_products')
PRICE_PATH = os.path.join(SOFA_DIR, '沙发经典产品价格表.xlsx')
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ==================== 1. 读取价格表：组合规格 + 价格 ====================
print(">>> 读取价格表组合规格...")
pythoncom.CoInitialize()
excel = win32.Dispatch('Excel.Application')
excel.Visible = False; excel.DisplayAlerts = False
wb = excel.Workbooks.Open(PRICE_PATH, 0, False, 5, '')

price_combos = {}  # {货号: [(规格名, 批发价)]}

for sheet_name in ['顾家经典固定', '顾家经典功能']:
    ws = wb.Sheets(sheet_name)
    data = ws.Range(ws.Cells(1,1), ws.Cells(ws.UsedRange.Rows.Count, ws.UsedRange.Columns.Count)).Value
    current_code = ''
    for i, row in enumerate(data):
        if i < 3: continue
        c2 = str(row[1]).strip() if row[1] is not None else ''
        c6 = str(row[5]).strip() if row[5] is not None else ''
        c7 = row[6] if row[6] is not None else 0
        if c2 and c2 != '货号' and re.match(r'^[\w\.]+$', c2):
            current_code = c2
        if current_code and c6 and c7:
            try: price = float(c7)
            except: price = 0
            # Only keep combined specs (with '+')
            if '+' in c6:
                price_combos.setdefault(current_code, []).append((c6.strip(), price))

wb.Close(); excel.Quit(); pythoncom.CoUninitialize()
print(f"  加载 {sum(len(v) for v in price_combos.values())} 个组合规格")

# ==================== 2. PPT文本读取 ====================
def get_pptx_texts(p):
    t = []
    try:
        for s in Presentation(p).slides:
            for sh in s.shapes:
                if sh.has_text_frame:
                    for pa in sh.text_frame.paragraphs:
                        x = pa.text.strip()
                        if x: t.append(x)
    except: pass
    return t

def get_pdf_texts(p):
    t = []
    try:
        for page in PdfReader(p).pages:
            x = page.extract_text()
            if x: t.extend([l.strip() for l in x.split('\n') if l.strip()])
    except: pass
    return t

def read_ppt(fp):
    ppt_dir = os.path.join(fp, 'PPT')
    if not os.path.isdir(ppt_dir): return []
    fs = os.listdir(ppt_dir)
    xlsx = [f for f in fs if f.endswith('.xlsx') and not f.startswith('~$')]
    pptx = [f for f in fs if f.endswith('.pptx')]
    pdf = [f for f in fs if f.endswith('.pdf')]
    if pptx: return get_pptx_texts(os.path.join(ppt_dir, pptx[0]))
    if pdf:  return get_pdf_texts(os.path.join(ppt_dir, pdf[0]))
    if xlsx: 
        try:
            wb = openpyxl.load_workbook(os.path.join(ppt_dir, xlsx[0]), data_only=True)
            ws = wb.active
            lines = []
            for r in range(1, ws.max_row+1):
                for c in range(1, ws.max_column+1):
                    v = ws.cell(r,c).value
                    if v: lines.append(str(v).strip())
            wb.close()
            return lines
        except: pass
    return []

# ==================== 3. 从PPT查找尺寸 ====================
def find_size_in_ppt(texts, spec_name, product_code=''):
    """在PPT文本中查找与规格匹配的尺寸"""
    # Product-specific fallback sizes (PPT name convention vs price table naming)
    FALLBACK_SIZES = {
        'JD.0036': {s[0]: s[1] for s in [
            ('3.5左+1.5右', 311), ('3左+1.5右A', 287), ('2.5左+1右', 251),
        ]},
        'JD.0006': {s[0]: s[1] for s in [
            ('3单+1.5单', 300), ('3单A+1.5单A', 276),
        ]},
        'JD.0069': {s[0]: s[1] for s in [
            ('3左+1.5右', 264), ('2.5左+1右', 240),
        ]},
        'JD.0020': {s[0]: s[1] for s in [
            ('3左+躺右', 312), ('3左+1.5右', 312),
            ('3左A+1.5右A', 288), ('2左+2右', 262),
        ]},
        'JD.0021': {s[0]: s[1] for s in [
            ('3单+1.5单', None), ('3单A+1.5单A', None),
        ]},
        'JD.0072': {s[0]: s[1] for s in [
            ('3左+1.5右A+凳组', 400), ('3.5左+1.5右+凳组', 400),
        ]},
        'JD.0077': {s[0]: s[1] for s in [
            ('3左+1.5右', 300), ('3左A+1.5右A', 270),
        ]},
    }
    
    # Check product-specific fallback first
    if product_code in FALLBACK_SIZES:
        for price_spec, size in FALLBACK_SIZES[product_code].items():
            if spec_name == price_spec or spec_name.startswith(price_spec + '+'):
                if size: return size
    
    # Direct match: look for spec_name followed by -size in PPT text
    spec_escaped = re.escape(spec_name)
    for t in texts:
        m = re.search(rf'{spec_escaped}[\s\-—]*(\d+(?:\.\d+)?)\s*(?:CM|cm|m)\b', t)
        if m:
            raw = float(m.group(1))
            return int(raw * 100) if raw < 10 else int(raw)
    
    # Extract all PPT spec→size mappings
    ppt_spec_sizes = {}
    for t in texts:
        # "L-大3双：311cm" → name="大3双" size=311
        ms = re.findall(r'[A-Z]+[\s\-—]*[：:]?\s*([\u4e00-\u9fff\w]+?)[\s\-—]*[：:]?\s*(\d{3,})\s*(?:CM|cm)\b', t)
        for name, size in ms:
            if len(name) >= 2 and int(size) >= 100:
                ppt_spec_sizes[name] = int(size)
    
    # Also extract plain letter:size mappings "S：214cm" → collect all sizes
    letter_sizes = []
    for t in texts:
        for size in re.findall(r'[A-Z]+\s*[：:]\s*(\d{3,})\s*(?:CM|cm)\b', t):
            if int(size) >= 100: letter_sizes.append(int(size))
    for i, s in enumerate(letter_sizes):
        ppt_spec_sizes[f'_letter_{i}'] = s
    
    if not ppt_spec_sizes:
        return 0
    
    # Try to match spec_name parts to PPT descriptions
    spec_norm = spec_name.replace(' ', '').replace('|', '').replace('+', '')
    
    # Try direct char overlap
    for ppt_name, size in sorted(ppt_spec_sizes.items(), key=lambda x: -x[1]):
        ppt_norm = ppt_name.replace(' ', '')
        # Check if any number from spec appears in ppt_name and vice versa
        spec_nums = set(re.findall(r'\d+\.?\d*', spec_norm))
        ppt_nums = set(re.findall(r'\d+\.?\d*', ppt_norm))
        if spec_nums & ppt_nums:
            return size
        # Check Chinese char overlap
        overlap = sum(1 for c in spec_norm if c in '\u4e00-\u9fff' and c in ppt_norm)
        if overlap >= 1:
            return size
    
    return 0

# Also extract size from the spec name itself (如 "大3双2电动317cm" 包含在名称中)
def extract_size_from_spec(spec):
    m = re.search(r'[（(].*?(\d{3,})\s*(?:CM|cm)?[）)]', spec)
    if m: return int(m.group(1))
    m = re.search(r'(\d{3,})\s*(?:CM|cm)\b', spec)
    if m: return int(m.group(1))
    return 0

def simplify_spec_name(spec):
    """简化规格名称：去掉尺寸标注、缩略显示"""
    # Remove parenthesized size info
    s = re.sub(r'[（(][^）)]*(\d{3,}\s*(?:CM|cm)?)[）)]', '', spec).strip()
    # Remove duplicate whitespace
    s = re.sub(r'\s+', ' ', s).strip()
    # Remove trailing | or +
    s = s.rstrip('|+ ')
    return s if s else spec

# ==================== 4. 提取PPT产品信息 ====================
def extract_ppt_info(texts):
    info = {'name': '', 'intro': '', 'features': [], 'fabric': '', 'filling': '', 'colors': []}
    # 产品名称
    for t in texts:
        t = t.strip()
        if t and len(t) > 2 and len(t) < 30:
            if not any(kw in t for kw in ['PPT', '核心', '魅力', '尺寸', '规格', '202', 'Pursuing']):
                if not info['name']: info['name'] = t
    # 面料
    for t in texts:
        if any(kw in t.lower() for kw in ['面料', '牛皮', '真皮', '仿皮', '布艺', '材质']):
            if not info['fabric']: info['fabric'] = re.split(r'[。\n]', t)[0][:200]
    # 填充
    for t in texts:
        if any(kw in t for kw in ['填充', '海绵', '羽绒', '羽毛']):
            if not info['filling']: info['filling'] = re.split(r'[。\n]', t)[0][:200]
    # 颜色/配色
    for t in texts:
        if any(kw in t for kw in ['配色', '颜色', '色彩']) and len(t) < 200:
            info['colors'].append(t[:200])
    # 特点
    for t in texts:
        if any(kw in t for kw in ['卖点', '魅力', '设计', '特点', '核心']) and len(t) > 10 and len(t) < 200:
            info['features'].append(t[:200])
    return info

MATERIAL_MAP = {'T':'真皮','W':'真皮','U':'真皮','O':'真皮','F':'仿皮','A':'布艺','C':'布艺','H':'布艺'}

def extract_colors_from_images(fp):
    """从浏览图文件名提取配色信息"""
    browse_dir = os.path.join(fp, '浏览图')
    if not os.path.isdir(browse_dir): return []
    colors = {}
    for fname in os.listdir(browse_dir):
        m = re.search(r'([A-Z]\d{5,}[\w\-]*)-([\u4e00-\u9fff]+)', fname)
        if m:
            code = m.group(1)
            name = m.group(2)
            first = code[0]
            material = MATERIAL_MAP.get(first, '未知')
            if code not in colors:
                colors[code] = (name, material)
    result = []
    for code, (name, material) in sorted(colors.items()):
        result.append(f"{code}-{name}（{material}）")
    return result

def extract_colors_from_ppt(texts):
    """从PPT文本提取配色信息"""
    result = []
    # Known color names for matching
    color_names = ['黑骑士','菊蕊白','暖栗色','暖沙','烟灰色','奶咖','栗子棕','魅影黑',
                   '星爵黑','岩灰','星云灰','琥珀橙','燕麦拿铁','砂岩棕','熔岩褐','烟墨棕',
                   '沙滩白','暗夜棕','芝士米','暮云灰','米白','暖灰','奶棕色','月光灰',
                   '奶油米','烟灰紫','浆果红棕','乳酪白','晨雾绿','烟粉','云石白','暮山褐',
                   '雾霾蓝','抹茶绿','宁和蓝','雾蓝','奶咖色','暖沙色','杏仁咖','可可棕',
                   '墨黛紫','云雀棕','栀子白','格调黑','摩卡棕','礁岩褐','贝母白','焦糖棕',
                   '枫糖棕','栗壳棕','烟熏橡木','黑曜石','醇巧黑','暖山玉']
    
    for t in texts:
        # Format 1: "U665010-星爵黑" (code-color_name)
        for code, name in re.findall(r'([A-Z]\d{5,}[\w\-]*)-([\u4e00-\u9fff]+)', t):
            first = code[0]
            material = MATERIAL_MAP.get(first, '未知')
            entry = f"{code}-{name}（{material}）"
            if entry not in result: result.append(entry)
        
        # Format 2: "T250061/F250061-X 菊蕊白（坚韧耐磨皮）" (code space color_name)
        for code, name in re.findall(r'([A-Z]\d{5,}(?:/[A-Z]\d{5,})*(?:-X)?)\s+([\u4e00-\u9fff/]+)(?=[（(]|\s|$)', t):
            name_clean = name.split('/')[0].strip()
            if name_clean and len(name_clean) <= 6:
                code_clean = code.split('/')[0].strip('-X')
                first = code_clean[0]
                material = MATERIAL_MAP.get(first, '未知')
                entry = f"{code_clean}-{name_clean}（{material}）"
                if entry not in result: result.append(entry)
        
        # Format 3: "可可棕T250122/F250122-X" (color_name directly followed by code)
        m = re.search(r'([\u4e00-\u9fff]{2,4})([A-Z]\d{5,}(?:/[A-Z]\d{5,})*(?:-X)?)', t)
        if m:
            name = m.group(1)
            code = m.group(2).split('/')[0].strip('-X')
            first = code[0]
            material = MATERIAL_MAP.get(first, '未知')
            entry = f"{code}-{name}（{material}）"
            if entry not in result: result.append(entry)
        
        # Format 4: "柔光肤感皮-U665162" (text-code) - only capture if text looks like color name
        for name, code in re.findall(r'([\u4e00-\u9fff/]{2,10})[-/]([A-Z]\d{5,}[\w\-]*)', t):
            name_clean = name.split('/')[0].strip('-皮料 ')
            code_clean = code.strip('-X')
            if name_clean and len(name_clean) <= 6:
                # Only keep if it's a known color name, not a material/process description
                if name_clean in color_names:
                    first = code_clean[0]
                    material = MATERIAL_MAP.get(first, '未知')
                    entry = f"{code_clean}-{name_clean}（{material}）"
                    if entry not in result: result.append(entry)
        
        # Format 5: "颜色有杏仁咖和星爵黑两个配色" - extract color names from context words
        if '配色' in t or '颜色' in t or ('色' in t and len(t) < 100):
            for cn in color_names:
                if cn in t:
                    # Try to find a nearby code
                    nearby = t[t.find(cn)-30:t.find(cn)+30]
                    code_m = re.search(r'([A-Z]\d{5,}[\w\-]*)', nearby)
                    if code_m:
                        code = code_m.group(1)
                        first = code[0]
                        material = MATERIAL_MAP.get(first, '未知')
                    else:
                        code = '-'
                        material = '未知'
                    entry = f"{code}-{cn}（{material}）" if code != '-' else f"颜色: {cn}"
                    if entry not in result: result.append(entry)
    
    # Format 6: Match standalone color names (entire line is just a color name)
    for t in texts:
        t_stripped = t.strip().strip('•►-— 0123456789')
        if t_stripped in color_names:
            entry = f"颜色: {t_stripped}"
            if entry not in result: result.append(entry)
    
    # Post-process: try to match codes to color names across separate lines
    # Collect all codes and color names
    codes_found = {}
    color_only = []
    for entry in result:
        m = re.match(r'([A-Z]\d{5,})-([\u4e00-\u9fff]+)（', entry)
        if m:
            codes_found[m.group(2)] = m.group(1)
    for entry in result:
        if entry.startswith('颜色: '):
            color_only.append(entry[4:])
    
    # If we have codes without matching colors, try to find from color_only
    if color_only:
        all_text = '\n'.join(texts)
        for cn in color_only:
            if cn not in codes_found:
                # Look for code near this color name
                idx = all_text.find(cn)
                if idx >= 0:
                    nearby = all_text[max(0,idx-50):idx+50]
                    cm = re.search(r'([A-Z]\d{5,})', nearby)
                    if cm:
                        code = cm.group(1)
                        first = code[0]
                        material = MATERIAL_MAP.get(first, '未知')
                        # Replace or add
                        new_entry = f"{code}-{cn}（{material}）"
                        result = [new_entry if e == f'颜色: {cn}' else e for e in result]
                        if new_entry not in result:
                            result.append(new_entry)
    
    return result

# ==================== 5. 工具函数 ====================
def safe_name(f): return f.replace(' ', '_').replace('/', '_')
def write_md(p, l):
    with open(p, 'w', encoding='utf-8') as f: f.write('\n'.join(l))
def price_str(v):
    return f"¥{v:,}" if v > 0 else ''

def list_images(d):
    if not os.path.isdir(d): return []
    exts = ('.jpg','.jpeg','.png','.gif','.webp','.bmp')
    return sorted([f for f in os.listdir(d) if f.lower().endswith(exts)])

def get_images(fp):
    r = {}
    for s in ['场景图','浏览图','入户实景图','白底图']:
        imgs = list_images(os.path.join(fp, s))
        if imgs: r[s] = imgs
    return r

def fmt_imgs(d, fp):
    lines = []
    for s, imgs in sorted(d.items()):
        lines.append(f"\n### {s} ({len(imgs)}张)\n")
        for img in imgs[:10]:
            rel = os.path.relpath(os.path.join(fp, s, img), BASE_DIR).replace('\\','/')
            lines.append(f"![{os.path.splitext(img)[0]}]({rel})")
        if len(imgs) > 10: lines.append(f"... 共{len(imgs)}张")
    return lines

# ==================== 6. 主流程 ====================
print(">>> 更新沙发 MD")
count = 0
for folder in sorted(os.listdir(SOFA_DIR)):
    fp = os.path.join(SOFA_DIR, folder)
    if not os.path.isdir(fp) or folder.startswith('~$'): continue
    
    texts = read_ppt(fp)
    ppt_info = extract_ppt_info(texts)
    images = get_images(fp)
    
    lines = [f"# {folder}\n"]
    
    # --- 产品信息 ---
    lines.append("## 产品信息\n")
    if ppt_info['name']:
        lines.append(f"**产品名称**: {ppt_info['name']}\n")
    lines.append(f"**产品型号**: {folder}\n")
    if ppt_info['intro']:
        lines.append(f"### 产品介绍\n{ppt_info['intro']}\n")
    if ppt_info['features']:
        lines.append("### 产品特点\n")
        for f in ppt_info['features'][:5]:
            lines.append(f"- {f}")
        lines.append("")
    if ppt_info['fabric']:
        lines.append(f"### 面料\n- {ppt_info['fabric']}\n")
    if ppt_info['filling']:
        lines.append(f"### 填充材质\n- {ppt_info['filling']}\n")
    if ppt_info['colors']:
        lines.append("### 配色方案\n")
        for c in ppt_info['colors'][:5]: lines.append(f"- {c}")
        lines.append("")
    
    # --- 配色信息（从浏览图+PPT提取） ---
    img_colors = extract_colors_from_images(fp)
    ppt_colors = extract_colors_from_ppt(texts)
    all_colors = list(dict.fromkeys(img_colors + ppt_colors))  # dedup preserving order
    if all_colors:
        lines.append("### 配色与材质\n")
        lines.append("| 色号 | 颜色名 | 材质 |")
        lines.append("|------|--------|------|")
        for entry in all_colors:
            m = re.match(r'([A-Z]\d{5,}[\w\-]*)-([\u4e00-\u9fff/]+)（(.+)）', entry)
            if m:
                lines.append(f"| {m.group(1)} | {m.group(2)} | {m.group(3)} |")
            elif entry.startswith('颜色: '):
                lines.append(f"| - | {entry[4:]} | - |")
            else:
                lines.append(f"| - | {entry} | - |")
        lines.append("")
    
    # --- 规格与价格（从价格表取组合规格） ---
    combos = price_combos.get(folder, [])
    # Also try matching by alternate code (e.g. JD.0061B → JD.0061)
    if not combos and '.' in folder:
        base = folder.split('.')[0] + '.' + folder.split('.')[1][0]
        combos = price_combos.get(base, [])
    
    if combos:
        lines.append("## 规格与价格\n")
        lines.append("| 规格 | 尺寸(CM) | 实际成交价 |")
        lines.append("|------|---------|-----------|")
        seen = set()
        for spec, ws_price in combos:
            actual = round(ws_price * 1.7)
            
            # Try to find size
            size = extract_size_from_spec(spec)
            if not size:
                size = find_size_in_ppt(texts, spec, folder)
            
            simple = simplify_spec_name(spec)
            
            # Dedup by simplified name
            key = (simple, size)
            if key in seen: continue
            seen.add(key)
            
            size_str = f"{size}cm" if size else '-'
            lines.append(f"| {simple} | {size_str} | {price_str(actual)} |")
        lines.append("")
    
    # --- 文档参考 ---
    ppt_dir = os.path.join(fp, 'PPT')
    if os.path.isdir(ppt_dir):
        for fname in os.listdir(ppt_dir):
            if fname.endswith(('.pptx','.pdf','.xlsx')) and not fname.startswith('~$'):
                lines.append(f"**参考文档**: {fname}\n")
                break
    
    # --- 图片 ---
    if images:
        lines.append("## 图片素材\n")
        lines.extend(fmt_imgs(images, fp))
        lines.append("")
    
    if not images and not combos:
        lines.append("*(暂无资料)*\n")
    
    write_md(os.path.join(OUTPUT_DIR, f"{safe_name(folder)}.md"), lines)
    count += 1
    ic = sum(len(v) for v in images.values())
    print(f"  {folder}: {len(combos)}组合 {ic}图")

print(f"\n完成: {count} 个沙发 MD")
