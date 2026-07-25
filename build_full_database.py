import os
import glob
import json
import re
import time
import pandas as pd
from pptx import Presentation
from openai import OpenAI
from typing import Dict

# ==================== 1. 核心配置区域 ====================
# API Key 配置（当前使用千问 DashScope API）
API_KEY = "sk-ws-H.EHLRIYY.lQ0t.MEUCIQDDAH_SsWSrRmqiMo2G8Ahp7wXYpoRc7Mj8Hw6Cj9FAdgIgMV-LYKqyrYAsGUfQqVIRHD6cXHqF0JTYMg50fD0cAkA"  # 👈 请替换为你的千问 API Key（DashScope）
BASE_URL = "https://dashscope.aliyuncs.com/compatible-mode/v1"  # 千问 DashScope API
MODEL_NAME = "qwen3.7-plus"     # 千问模型（使用更强模型以获得更精准的档案生成）

# 本地文件路径设置（支持绝对路径和相对路径，统一使用 Unix 风格）
BASE_DATA_DIR = "./sofa_products"         # 包含几十个沙发产品子文件夹的主目录
PRICE_EXCEL_PATH = "./sofa_products/沙发经典产品价格表.xlsx"          # 价格表 Excel 文件
OUTPUT_MD_DIR = "./markdown_db"           # 生成的 Markdown 档案保存目录
JSON_OUTPUT_FILE = "product_images.json"  # 生成的图片索引 JSON 文件

# 初始化 AI 客户端
client = OpenAI(api_key=API_KEY, base_url=BASE_URL)


# ==================== 2. 数据处理辅助函数 ====================
def read_excel_prices(excel_path: str) -> tuple:
    """
    读取 Excel 价格表（支持多 sheet），忽略零售价，按内部定价规则计算【实际成交价】

    只保留组合型规格（规格中含 '+'），单品规格不报价。
    
    返回: (price_dict, module_prices)
      - price_dict: { 型号大写: [{"型号": str, "规格": str, "实际成交价": str, "批发价": float}] }
      - module_prices: { 型号大写: [{"型号": str, "规格": str, "实际成交价": str}] }
    """
    if not os.path.exists(excel_path):
        print(f"⚠️ 提示：未找到价格表【{excel_path}】，本次将不关联价格信息。")
        return {}, {}

    try:
        xls = pd.ExcelFile(excel_path)
        skip_sheets = {'WpsReserved_CellImgList'}
        target_sheets = [s for s in xls.sheet_names if s not in skip_sheets]

        price_dict: Dict[str, list] = {}
        module_prices: Dict[str, list] = {}

        for sheet_name in target_sheets:
            df = pd.read_excel(xls, sheet_name=sheet_name, header=1)

            # 自动识别型号列、批发价列和规格列
            model_col = _find_column(df, ['型号', 'model', '产品型号', '货号', '编号'])
            price_col = _find_column(df, ['批发价', 'wholesale', '出厂价', '成本价'])
            spec_col = _find_column(df, ['规格', 'spec', '尺寸规格'])

            if not model_col or not price_col or not spec_col:
                print(f"  ⏩ Sheet【{sheet_name}】未找到完整列，已跳过（表头: {list(df.columns)}）")
                continue

            # 记录上一个有型号的行（用于组合行没有型号时回填关联）
            last_model = None

            for idx, row in df.iterrows():
                # 跳过子表头行
                if idx == 0:
                    continue

                model_val = row[model_col]
                spec_val = str(row[spec_col]).strip() if spec_col and pd.notna(row[spec_col]) else ''

                # 跳过二级表头行
                if pd.notna(model_val) and str(model_val).strip() == '货号':
                    continue

                # 判断是否是组合型规格（包含 '+'）
                is_combo = '+' in spec_val

                # 如果有型号，更新 last_model
                if pd.notna(model_val):
                    model = str(model_val).strip().upper()
                    if model and model != 'NAN':
                        last_model = model
                else:
                    # 没有型号：如果是组合规格，尝试关联到上一个产品
                    if not is_combo:
                        continue
                    model = last_model

                if not model:
                    continue

                raw_price = row[price_col]
                try:
                    wholesale_price = float(raw_price)
                    actual_price = int(round(wholesale_price * 1.7))
                    formatted_price = f"¥{actual_price:,}"
                except (ValueError, TypeError):
                    formatted_price = "暂无定价"

                if is_combo:
                    if model not in price_dict:
                        price_dict[model] = []
                    price_dict[model].append({
                        "型号": model,
                        "规格": spec_val,
                        "实际成交价": formatted_price,
                        "批发价": float(raw_price) if isinstance(raw_price, (int, float)) else 0,
                    })
                else:
                    if model not in module_prices:
                        module_prices[model] = []
                    module_prices[model].append({
                        "型号": model,
                        "规格": spec_val,
                        "实际成交价": formatted_price,
                    })

            print(f"  ✅ Sheet【{sheet_name}】处理完成")

        combo_count = sum(len(v) for v in price_dict.values())
        module_count = sum(len(v) for v in module_prices.values()) if module_prices else 0
        print(f"📊 共加载 {combo_count} 条组合型 + {module_count} 条模块单价价格数据，已全自动换算【实际成交价】！")
        return price_dict, module_prices
    except Exception as e:
        print(f"❌ 读取 Excel 价格表失败: {e}")
        return {}, {}


def _find_column(df: pd.DataFrame, candidates: list) -> str | None:
    """在 DataFrame 中根据候选名称列表匹配列名（忽略大小写）"""
    col_lower = {c.strip().lower(): c for c in df.columns}
    for candidate in candidates:
        key = candidate.strip().lower()
        if key in col_lower:
            return col_lower[key]
    return None


def build_price_lookup(price_db: Dict[str, list]) -> Dict[str, list]:
    """
    构建型号->组合价格的快速查找表。
    将原始 price_db key（型号全称）展开为按空格/斜杠分隔的子 token，
    使后续文件夹名匹配更灵活。
    """
    lookup: Dict[str, list] = {}
    for model_key, info_list in price_db.items():
        lookup[model_key] = info_list
        for token in re.split(r'[\s/_-]+', model_key):
            token = token.strip()
            if len(token) >= 2:
                lookup[token] = info_list
    return lookup


def extract_text_from_ppt(ppt_path: str) -> str:
    """从 PPT 文件中提取所有文本与卖点关键词"""
    try:
        prs = Presentation(ppt_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            full_text.append(text)
        return "\n".join(full_text)
    except Exception as e:
        print(f"  ⚠️ 读取 PPT 失败 ({os.path.basename(ppt_path)}): {e}")
        return ""


def extract_text_from_pdf(pdf_path: str) -> str:
    """从 PDF 文件中提取所有文本"""
    try:
        import fitz
        doc = fitz.open(pdf_path)
        full_text = []
        for page in doc:
            full_text.append(page.get_text())
        doc.close()
        return "\n".join(full_text)
    except Exception as e:
        print(f"  ⚠️ 读取 PDF 失败 ({os.path.basename(pdf_path)}): {e}")

def extract_text_from_docx(docx_path: str) -> str:
    """从 Word 文件中提取所有文本"""
    try:
        from docx import Document
        doc = Document(docx_path)
        full_text = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(full_text)
    except Exception as e:
        print(f"  \u26a0\ufe0f 读取 Word 失败 ({os.path.basename(docx_path)}): {e}")
        return ""


def extract_text_from_xlsx(xlsx_path: str) -> str:
    """从 Excel 文件中提取所有文本内容"""
    try:
        import pandas as pd
        xls = pd.ExcelFile(xlsx_path)
        all_text = []
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet)
            texts = []
            for col in df.columns:
                for val in df[col].dropna():
                    s = str(val).strip()
                    if s and s != 'nan':
                        texts.append(s)
            if texts:
                all_text.append(f"[Sheet: {sheet}]\n" + "\n".join(texts[:200]))
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"  \u26a0\ufe0f 读取 Excel 失败 ({os.path.basename(xlsx_path)}): {e}")
        return ""


def extract_text_from_image(image_path: str) -> str:
    """从图片文件中通过千问多模态 API 提取文本（识别详情页的面料、填充等信息）"""
    import base64
    try:
        with open(image_path, "rb") as f:
            img_b64 = base64.b64encode(f.read()).decode("utf-8")

        ext = os.path.splitext(image_path)[1].lower().replace(".", "")
        mime = "png" if ext == "png" else "jpeg"

        response = client.chat.completions.create(
            model="qwen-vl-plus",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/{mime};base64,{img_b64}"}},
                        {"type": "text", "text": "请完整提取这张产品详情页中的所有文字信息，特别是：产品名称、面料材质、填充物类型、海绵密度、尺寸规格、颜色等参数。按原文输出，不要遗漏。"},
                    ],
                }
            ],
            temperature=0.1,
            max_tokens=2000,
        )
        return response.choices[0].message.content or ""
    except Exception as e:
        print(f"  \u26a0\ufe0f OCR 识别图片失败 ({os.path.basename(image_path)}): {e}")
        return ""


def scan_images_in_folder(folder_path: str, folder_name: str) -> dict:
    """扫描文件夹中的图片，按【浏览图】【场景图】分类，并从文件名提取配色关联信息"""
    image_data = {
        "catalog_images": [],    # 浏览图/规格图
        "scene_images": [],      # 实拍场景图
        "color_images": {},      # { 色号或配色名: {"code": str, "name": str, "files": [路径...]} }
    }

    all_imgs = (
        glob.glob(os.path.join(folder_path, "**", "*.[jJ][pP][gG]"), recursive=True)
        + glob.glob(os.path.join(folder_path, "**", "*.[pP][nN][gG]"), recursive=True)
        + glob.glob(os.path.join(folder_path, "**", "*.[jJ][pP][eE][gG]"), recursive=True)
    )

    for img_path in all_imgs:
        rel_path = os.path.relpath(img_path, start=".").replace("\\", "/")
        filename = os.path.basename(img_path)

        # 使用标准化目录名分类
        if "/场景图/" in rel_path or "/入户实景图/" in rel_path:
            image_data["scene_images"].append(rel_path)
        elif "/浏览图/" in rel_path:
            image_data["catalog_images"].append(rel_path)
            # 从文件名提取颜色信息
            color_info = _extract_color_from_filename(filename, folder_name)
            if color_info:
                cid = color_info.get("code") or color_info.get("name", "")
                prefix = color_info.get("prefix", folder_name)
                unique_key = f"{prefix}|{cid}"
                if unique_key not in image_data["color_images"]:
                    color_info["files"] = []
                    color_info["prefix"] = prefix
                    image_data["color_images"][unique_key] = color_info
                image_data["color_images"][unique_key]["files"].append(rel_path)

    return image_data


def _extract_color_from_filename(filename: str, folder_name: str) -> dict | None:
    """从文件名提取颜色信息，支持多种命名格式"""
    name_no_ext = os.path.splitext(filename)[0]
    escaped = re.escape(folder_name)

    # --- 精确前缀匹配模式 ---
    # 模式1: {型号}-{色号}-{配色名}（已重命名的文件）
    m = re.match(rf'^{escaped}[-_]([A-Z]\d+)[-_]([\u4e00-\u9fff]+.*)$', name_no_ext)
    if m:
        return {"code": m.group(1), "name": m.group(2).strip(), "prefix": folder_name}

    # 模式2: {型号}-{色号}
    m = re.match(rf'^{escaped}[-_]([A-Z]\d+)$', name_no_ext)
    if m:
        return {"code": m.group(1), "name": None, "prefix": folder_name}

    # 模式3: {型号}-{中文配色名}（可能含后缀如"大转角"）
    m = re.match(rf'^{escaped}[-_]([\u4e00-\u9fff]+.*)$', name_no_ext)
    if m:
        return {"code": None, "name": m.group(1).strip(), "prefix": folder_name}

    # 模式4: {型号}空格{配色名}空格{色号}-{额外文字}
    m = re.match(rf'^{escaped}\s+([\u4e00-\u9fff]+)\s+([A-Z]\d+).*$', name_no_ext)
    if m:
        return {"code": m.group(2), "name": m.group(1).strip(), "prefix": folder_name}

    # 模式5: {型号}空格{配色名}
    m = re.match(rf'^{escaped}\s+([\u4e00-\u9fff]+.*)$', name_no_ext)
    if m:
        return {"code": None, "name": m.group(1).strip(), "prefix": folder_name}

    # --- 变体前缀匹配（如 JD.6013B, JD0062 等） ---
    # 提取文件名的实际前缀（以型号格式开头的部分）
    actual_prefix_m = re.match(r'^([A-Za-z0-9.]+)', name_no_ext)
    if actual_prefix_m:
        actual_prefix = actual_prefix_m.group(1)
        # 检查实际前缀是否以文件夹名为基础扩展（如 JD.6013 → JD.6013B）
        if actual_prefix.startswith(folder_name) and actual_prefix != folder_name:
            var_escaped = re.escape(actual_prefix)
            # 模式6: {变体前缀}空格{配色名}空格{色号}
            m = re.match(rf'^{var_escaped}\s+([\u4e00-\u9fff]+)\s+([A-Z]\d+).*$', name_no_ext)
            if m:
                return {"code": m.group(2), "name": m.group(1).strip(), "prefix": actual_prefix}
            # 模式7: {变体前缀}空格{配色名}
            m = re.match(rf'^{var_escaped}\s+([\u4e00-\u9fff]+.*)$', name_no_ext)
            if m:
                return {"code": None, "name": m.group(1).strip(), "prefix": actual_prefix}
            # 模式8: {变体前缀}-{色号}
            m = re.match(rf'^{var_escaped}[-_]([A-Z]\d+)$', name_no_ext)
            if m:
                return {"code": m.group(1), "name": None, "prefix": actual_prefix}

    # --- 前缀模糊匹配（处理无点变体如 JD0062 → 匹配 JD.0062） ---
    fuzzy_prefix = re.escape(folder_name.replace('.', ''))
    m = re.match(rf'^{fuzzy_prefix}\s+([\u4e00-\u9fff]+)\s+([A-Z]\d+).*$', name_no_ext)
    if m:
        return {"code": m.group(2), "name": m.group(1).strip(), "prefix": folder_name}

    # 模式9: {型号}-{色号}_{色号2}-X（双色号）
    m = re.match(rf'^{escaped}[-_]([A-Z]\d+)_.*$', name_no_ext)
    if m:
        return {"code": m.group(1), "name": None, "prefix": folder_name}

    return None


def enrich_color_names_from_markdown(md_content: str, color_images: dict) -> dict:
    """
    从 AI 生成的 Markdown 中解析配色映射（色号↔配色名），双向回填到 color_images
    支持多种 Markdown 格式：
      - "U665010 星爵黑" / "U665010-星爵黑"（同行）
      - "**雾咖色**...\n  - 真皮材质：T250032"（跨行配色块）
    """
    if not color_images:
        return color_images

    code_to_name = {}
    name_to_code = {}

    # 模式A: 跨行配色块（JD.6025风格）
    # - **雾咖色**（说明文字）
    #   - 真皮材质：T250032
    block_pattern = re.compile(
        r'\*\*([\u4e00-\u9fff]+)色?\*\*[^\n]*\n\s*[-–—*+]+\s*[^：:\n]*[：:]\s*([A-Z]\d+)',
        re.MULTILINE,
    )
    for m in block_pattern.finditer(md_content):
        name_raw, code = m.group(1).strip(), m.group(2)
        code_to_name[code] = name_raw
        name_to_code[name_raw] = code

    # 模式B: 同行 "色号 配色名"（空格/短横/冒号分隔，不跨行）
    inline_pattern = re.compile(r'([A-Z]\d{3,})[ \t]*[-—: ]+[ \t]*([\u4e00-\u9fff]+)')
    for m in inline_pattern.finditer(md_content):
        code, name_raw = m.group(1), m.group(2)
        code_to_name[code] = name_raw
        name_to_code[name_raw] = code

    # 模式C: **配色名色**：色号（配色名在前）
    # 如: **雾咖色**：T250032
    name_first = re.compile(r'\*\*([\u4e00-\u9fff]+)色?\*\*\s*[：:]\s*([A-Z]\d+)')
    for m in name_first.finditer(md_content):
        name_raw, code = m.group(1).strip(), m.group(2)
        code_to_name[code] = name_raw
        name_to_code[name_raw] = code

    # 回填：已有 code 的补 name，已有 name 的补 code，同时添加反向条目
    new_entries = {}
    for cid, info in list(color_images.items()):
        code = info.get("code")
        name = info.get("name")

        # code → name
        if code and code in code_to_name:
            info["name"] = code_to_name[code]

        # name → code
        if name and name in name_to_code:
            info["code"] = name_to_code[name]

        # 如果有 code 但没有 name 条目，新增 name 作为 key 的反向引用
        if code and code in code_to_name:
            name_key = code_to_name[code]
            if name_key not in color_images and name_key not in new_entries:
                new_entries[name_key] = {"code": code, "name": name_key, "files": list(info["files"])}

        # 如果有 name 但没有 code 条目，新增 code 作为 key 的反向引用
        if name and name in name_to_code:
            code_key = name_to_code[name]
            if code_key not in color_images and code_key not in new_entries:
                new_entries[code_key] = {"code": code_key, "name": name, "files": list(info["files"])}

    color_images.update(new_entries)
    return color_images


def rename_color_image_files(color_images: dict, catalog_images: list, scene_images: list) -> None:
    """
    将配色图片统一重命名，优先使用 {前缀}-{色号}-{配色名}.ext 格式，
    如无色号则用 {前缀}-{配色名}.ext，如无配色名则用 {前缀}-{色号}.ext
    前缀从文件名提取的实际型号前缀（如 JD.0061 / JD.6013B）
    """
    old_to_new = {}

    for cid, info in color_images.items():
        code = info.get("code")
        name = info.get("name")
        prefix = info.get("prefix", "")
        if not prefix:
            continue
        if not code and not name:
            continue

        old_files = info.get("files", [])
        new_files = []
        for old_rel in old_files:
            dir_name = os.path.dirname(old_rel)
            ext = os.path.splitext(old_rel)[1]
            # 构造新文件名：有code和name用完整格式，否则用可用的部分
            if code and name:
                new_basename = f"{prefix}-{code}-{name}{ext}"
            elif name:
                new_basename = f"{prefix}-{name}{ext}"
            else:
                new_basename = f"{prefix}-{code}{ext}"
            new_rel = os.path.join(dir_name, new_basename).replace("\\", "/")

            if old_rel != new_rel and os.path.exists(old_rel.replace("/", "\\")):
                old_abs = old_rel.replace("/", "\\")
                new_abs = new_rel.replace("/", "\\")
                if not os.path.exists(new_abs):
                    os.rename(old_abs, new_abs)
                    print(f"    📸 重命名: {os.path.basename(old_abs)} → {new_basename}")
            new_files.append(new_rel)
            old_to_new[old_rel] = new_rel
        info["files"] = new_files

    # 更新 catalog_images 和 scene_images 中的路径
    def _update_list(img_list):
        for i, p in enumerate(img_list):
            if p in old_to_new:
                img_list[i] = old_to_new[p]

    _update_list(catalog_images)
    _update_list(scene_images)


# ==================== 2.5 沙发类型检测 ====================
def detect_sofa_type(price_info_list: list, module_info_list: list, doc_text: str = "") -> str:
    """
    根据价格数据和文档（PPT/PDF）文本检测沙发类型：
    - 固定沙发：简单规格如 "1.5左"、"3右"
    - 功能沙发：含电动/功能字段，如 "1.5左电动|手靠|坐宽85"
    - 按摩沙发：含按摩功能字段
    """
    # 从价格数据中收集所有规格文本
    all_specs = []
    for p in (price_info_list or []):
        all_specs.append(p.get('规格', ''))
    for p in (module_info_list or []):
        all_specs.append(p.get('规格', ''))
    specs_text = ' '.join(all_specs)

    combined_text = specs_text + ' ' + doc_text

    # 检测优先级：按摩 > 电动功能 > 固定
    has_massage = '按摩' in combined_text
    has_power = '电动' in combined_text or '功能架' in combined_text or '零重力' in combined_text

    if has_massage:
        # 按摩沙发是功能沙发的子类
        return '功能沙发'
    if has_power:
        return '功能沙发'
    return '固定沙发'


def build_type_spec_guide(sofa_type: str) -> str:
    """根据沙发类型返回规格格式说明（按摩沙发是功能沙子的子类，共用同一套格式）"""
    guides = {
        '功能沙发': """
- 本产品是【功能沙发】（含按摩功能款式），规格由多个模块用"+"拼接，每个模块内用"|"分隔属性
- 典型格式如："1.5左电动|手靠|按摩功能|坐宽85|无电源+1.5无组|手靠|坐宽85+1.5右电动|手靠|坐宽85（大3双2电动313cm）"
- 或："1.5左电动|储物|无电源+1.5无组电动+1.5右电动|储物|无电源"
- 模块属性可能包含：电动/手动、手靠、坐宽XX、储物、无电源、USB、扶手翻折、按摩功能、隐藏杯托等
- 组合总长通常标注在末尾括号中，如"（大3双2电动313cm）"
- 请使用价格表中的真实规格名称，不要简化或编造""",
        '固定沙发': """
- 本产品是【固定沙发】，规格由多个模块用"+"拼接
- 简单格式如："3左+1.5右"、"2.5左+1右"
- 如有扶手翻折等功能属性，用"|"分隔，如"3.5左|扶手翻折+1.5右"
- 请使用价格表中的真实规格名称，不要简化或编造""",
    }
    return guides.get(sofa_type, guides['固定沙发'])


# ==================== 3. 核心主程序 ====================
def main():
    if not os.path.exists(OUTPUT_MD_DIR):
        os.makedirs(OUTPUT_MD_DIR)

    # 1. 载入 Excel 价格表数据并构建快速查找表
    price_db, module_db = read_excel_prices(PRICE_EXCEL_PATH)
    price_lookup = build_price_lookup(price_db) if price_db else {}
    all_product_images_json: Dict[str, dict] = {}

    # 2. 检查主数据文件夹
    if not os.path.exists(BASE_DATA_DIR):
        print(f"❌ 错误：找不到主文件夹【{BASE_DATA_DIR}】，请先建好文件夹并放入产品数据！")
        return

    subfolders = [f.path for f in os.scandir(BASE_DATA_DIR) if f.is_dir()]
    print(f"🔍 找到 {len(subfolders)} 个产品子文件夹，开始构建知识库...\n")

    for folder in subfolders:
        folder_name = os.path.basename(folder)
        print(f"📦 正在自动化归纳产品: 【{folder_name}】")

        # A. 搜寻并提取 PPT/PDF 文本内容（从标准化 PPT/ 目录读取）
        ppt_dir = os.path.join(folder, "PPT")
        doc_files = []
        if os.path.exists(ppt_dir):
            doc_files = (
                glob.glob(os.path.join(ppt_dir, "*.pptx"))
                + glob.glob(os.path.join(ppt_dir, "*.ppt"))
                + glob.glob(os.path.join(ppt_dir, "*.pdf"))
                + glob.glob(os.path.join(ppt_dir, "*.docx"))
                + glob.glob(os.path.join(ppt_dir, "*.xlsx"))
                + glob.glob(os.path.join(ppt_dir, "*.[jJ][pP][gG]"))
                + glob.glob(os.path.join(ppt_dir, "*.[jJ][pP][eE][gG]"))
                + glob.glob(os.path.join(ppt_dir, "*.[pP][nN][gG]"))
            )
        doc_files = [f for f in doc_files if not os.path.basename(f).startswith('~$')]

        doc_text = ""
        for doc in doc_files:
            ext = os.path.splitext(doc)[1].lower()
            if ext == ".pdf":
                doc_text += extract_text_from_pdf(doc) + "\n"
            elif ext == ".docx":
                doc_text += extract_text_from_docx(doc) + "\n"
            elif ext == ".xlsx":
                doc_text += extract_text_from_xlsx(doc) + "\n"
            elif ext in (".jpg", ".jpeg", ".png"):
                doc_text += extract_text_from_image(doc) + "\n"
            else:
                doc_text += extract_text_from_ppt(doc) + "\n"

        # B. 搜寻并分类图片路径
        images_info = scan_images_in_folder(folder, folder_name)
        all_product_images_json[folder_name] = images_info

        # C. 匹配组合型价格数据
        price_info_list = []
        folder_tokens = re.split(r'[\s/_-]+', folder_name.upper().strip())
        # 额外提取纯字母数字型号（如从"JD.0080图"提取"JD.0080"）
        folder_core_for_price = re.sub(r'[^a-zA-Z0-9.]+', '', folder_name.upper())
        if folder_core_for_price and folder_core_for_price not in folder_tokens:
            folder_tokens.append(folder_core_for_price)
        for token in folder_tokens:
            if token in price_lookup:
                price_info_list = price_lookup[token]
                break
        if not price_info_list and folder_name.upper() in price_lookup:
            price_info_list = price_lookup[folder_name.upper()]

        # C2. 匹配模块单价（用于类型检测和无组合价格时的参考）
        module_info_list = []
        if module_db:
            for token in folder_tokens:
                if token in module_db:
                    module_info_list = module_db[token]
                    break

        # C3. 检测沙发类型
        sofa_type = detect_sofa_type(price_info_list, module_info_list, doc_text)
        type_spec_guide = build_type_spec_guide(sofa_type)
        print(f"  🏷️ 类型识别: {sofa_type}")

        # D. 构造 AI Prompt
        system_prompt = "你是一位家居产品数据归档专家。请根据下方提供的PPT/PDF原文和价格表数据，提取真实的产品信息，生成结构化的Markdown档案。"

        # 配色材质编码说明
        material_code_guide = """【配色材质编码规则】：
配色编号的首字母代表材质类型，请在列出配色时标注材质：
  - T、W、U、O 开头 → 真皮
  - F 开头 → 仿皮
  - A、C、H 开头 → 布艺

【全真皮 vs 真皮区分】：
  - 如果该配色只有 T/W/U/O 开头（无 F 开头）→ 全真皮（所有面均为真皮）
  - 如果该配色同时有 T/W/U/O 和 F 开头 → 真皮（接触面真皮，非接触面仿皮）
示例：
  - U665010 → 全真皮（仅有U开头，无F开头）
  - T250032 / F250032-X → 真皮（接触面T开头真皮，非接触面F开头仿皮）"""

        # 构造价格描述文本
        has_combo_price = bool(price_info_list)
        if has_combo_price:
            price_lines = []
            for p in price_info_list:
                price_lines.append(f"  - 规格「{p['规格']}」→ 实际成交价 {p['实际成交价']}")
            price_desc = "以下是价格表中该产品的组合型规格及对应实际成交价，请以这些规格为尺寸规格的主要来源：\n" + "\n".join(price_lines)
        elif module_info_list:
            # 没有组合价但有模块单价，展示模块信息供参考
            module_lines = []
            for p in module_info_list[:15]:  # 展示前15条模块单价
                module_lines.append(f"  - 模块「{p['规格']}」→ 参考价 {p['实际成交价']}")
            price_desc = "（该产品在价格表中未找到组合型规格报价，仅找到以下模块单价供参考）\n" + "\n".join(module_lines)
        else:
            price_desc = "（该产品在价格表中无对应数据，当前仅根据产品文档生成基础框架）"

        has_doc = bool(doc_text.strip())
        doc_section = doc_text if has_doc else '（未提取到产品文档文本）'

        if has_doc:
            user_prompt = f"""文档原文：
{doc_section}

请从以上文档（PPT/PDF）原文中提取产品信息。

【重要要求】：
1. 【尺寸规格部分】请严格遵循以下格式指导：{type_spec_guide}
2. 如果价格表中已有该产品的规格数据，优先使用价格表中的真实规格名称，不要自行简化或编造
3. 核心卖点从文档中提炼，用简洁语言概括
4. 面料材质和配色从文档中提取真实信息
5. {material_code_guide}

{price_desc}

请严格按照以下Markdown结构输出：

# 产品完整品名与型号

## 1. 核心卖点与设计亮点
（从文档原文中提炼3-5个真实卖点，每个卖点用一句话概括）

## 2. 尺寸与规格参数
（请使用价格表中的真实组合规格，用Markdown表格展示：规格名称、尺寸/模块组成、备注说明）
（示例表格格式：第一列写完整规格如"3.5左|扶手翻折+1.5右|扶手翻折"，第二列写组合总长或模块详情）

## 3. 面料材质与配色选择
（从产品详情页/文档中完整提取面料材质、填充物信息，包括：接触面材质、填充海绵类型/密度、座包/靠包/扶手填充物等规格参数，以及配色信息）

## 4. 价格与市场定位
（基于价格表数据展示各规格及对应成交价，规格名称必须完整保留括号内的尺寸信息如"（大3双2电动313cm）"）
"""
        else:
            user_prompt = f"""请根据文件夹名称【{folder_name}】整理基础产品框架。

【重要要求】：
1. 类型识别：本产品是{sofa_type}
2. 尺寸规格格式指导：{type_spec_guide}
3. 如有价格数据，优先使用
4. {material_code_guide}

{price_desc}

请严格按照以下Markdown结构输出：

# 产品完整品名与型号

## 1. 核心卖点与设计亮点

## 2. 尺寸与规格参数

## 3. 面料材质与配色选择

## 4. 价格与市场定位
"""

        # E. 调用大模型提炼与生成（含自动重试，应对限流）
        md_content = None
        max_retries = 5
        for attempt in range(max_retries):
            try:
                response = client.chat.completions.create(
                    model=MODEL_NAME,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt},
                    ],
                    temperature=0.2,
                )
                md_content = response.choices[0].message.content
                break  # 成功则跳出重试循环
            except Exception as e:
                error_msg = str(e)
                if "429" in error_msg or "rate limit" in error_msg.lower() or "tpm" in error_msg.lower():
                    wait = 2 ** attempt  # 指数退避：1s, 2s, 4s, 8s, 16s
                    print(f"  ⏳ 遇到限流（第{attempt+1}次），等待{wait}秒后重试...")
                    time.sleep(wait)
                else:
                    print(f"  ❌ AI 提炼失败: {e}")
                    break  # 非限流错误，跳出循环

        if md_content:
            # 从 AI 生成的 Markdown 中解析配色名，回填到图片索引
            enrich_color_names_from_markdown(md_content, images_info.get("color_images", {}))

            # 统一重命名配色图片为 {前缀}-{色号}-{配色名}.ext 格式
            rename_color_image_files(
                images_info.get("color_images", {}),
                images_info.get("catalog_images", []),
                images_info.get("scene_images", []),
            )

            # 保存 Markdown 文件
            md_filename = f"{folder_name}.md"
            md_path = os.path.join(OUTPUT_MD_DIR, md_filename)
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(md_content)
            print(f"  ✅ Markdown 档案生成完毕: {md_path}")
        else:
            print(f"  ❌ AI 提炼失败: 无法生成 {folder_name} 的 Markdown")

    # 3. 输出全局图片 mapping 索引 JSON 文件（放到 OUTPUT_MD_DIR 内统一管理）
    json_path = os.path.join(OUTPUT_MD_DIR, JSON_OUTPUT_FILE)
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(all_product_images_json, f, ensure_ascii=False, indent=4)

    print("\n==========================================")
    print(f"🎉 全部数据归纳与换算完成！")
    print(f"📄 Markdown 产品档案库 ➡️ 【{OUTPUT_MD_DIR}/】")
    print(f"🖼️ 图文索引 JSON 数据库 ➡️ 【{json_path}】")
    print("==========================================")


if __name__ == "__main__":
    main()
