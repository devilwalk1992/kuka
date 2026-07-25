#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
提取沙发产品PPT中的适配风格和规格尺寸信息
"""
import os
import re
import sys
from pptx import Presentation
from pypdf import PdfReader

BASE_DIR = r"G:\Trae工作文件\顾家产品库\sofa_products"

# 要分析的文件列表: (产品ID, 文件名)
FILES = [
    ("9659", "9659.pptx"),
    ("JD.0006", "JD.0006.pptx"),
    ("JD.0036", "JD.0036.pptx"),
    ("JD.0061", "01 JD.0061大适界+A6003+PT3760系列+PT3323Y(-A) 客餐厅空间培训课件（260313）.pptx"),
    ("JD.0062", "02 JD.0062 大自在+PT3319系列 客餐厅空间培训课件（260520）.pptx"),
    ("HS.8002", "HS.8002.pdf"),
    ("JD.0020", "JD.0020.pdf"),
]

# 关键词 - 风格相关
STYLE_KEYWORDS = ["风格", "适配", "现代", "奶油", "意式", "轻奢", "北欧", "极简", "中式",
                  "复古", "工业", "侘寂", "法式", "美式", "日式", "简约", "田园", "新中式",
                  "欧式", "ins", "INS", "网红", "韩式"]

# 关键词 - 尺寸相关
DIM_KEYWORDS = ["尺寸", "规格", "CM", "cm", "厘米", "长", "宽", "高", "深",
                "×", "*", "x", "mm", "MM", "单人位", "双人位", "三人位",
                "转角", "脚踏", "贵妃", "躺位", "单位", "双位", "三位", "组合"]

# 关键词 - 沙发类型
SOFA_KEYWORDS = ["沙发", "单人", "双人", "三人", "转角", "贵妃", "躺位", "脚踏"]

OUTPUT_FILE = r"G:\Trae工作文件\顾家产品库\ppt_analysis_result.txt"


def log(msg=""):
    print(msg)
    with open(OUTPUT_FILE, "a", encoding="utf-8") as f:
        f.write(msg + "\n")


def extract_text_from_pptx(filepath):
    """从PPTX文件提取所有文本"""
    texts = []
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(row_texts))
            if slide_texts:
                texts.append((i, slide_texts))
    except Exception as e:
        texts.append((0, [f"[读取错误] {e}"]))
    return texts


def extract_text_from_pdf(filepath):
    """从PDF文件提取所有文本"""
    texts = []
    try:
        reader = PdfReader(filepath)
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                lines = [line.strip() for line in text.split("\n") if line.strip()]
                texts.append((i, lines))
    except Exception as e:
        texts.append((0, [f"[读取错误] {e}"]))
    return texts


def is_relevant_line(line):
    """判断一行文本是否需要关注"""
    combined_keywords = STYLE_KEYWORDS + DIM_KEYWORDS + SOFA_KEYWORDS
    return any(kw in line for kw in combined_keywords)


def analyze_product(product_id, filename):
    """分析一个产品的PPT，提取相关信息"""
    filepath = os.path.join(BASE_DIR, product_id, "PPT", filename)
    log(f"\n{'='*80}")
    log(f"📦 产品编号: {product_id}")
    log(f"📁 文件: {filename}")
    log(f"{'='*80}")

    if not os.path.exists(filepath):
        log(f"❌ 文件不存在: {filepath}")
        return

    # 提取文本
    if filename.lower().endswith(".pptx"):
        slides = extract_text_from_pptx(filepath)
    elif filename.lower().endswith(".pdf"):
        slides = extract_text_from_pdf(filepath)
    else:
        log(f"❌ 不支持的文件格式: {filename}")
        return

    # 收集所有文本
    all_lines = []
    for slide_num, lines in slides:
        for line in lines:
            all_lines.append(line)

    # ---------- 风格分析 ----------
    log(f"\n▶ 1. 适配风格分析")
    style_info = []
    for line in all_lines:
        found_styles = [kw for kw in STYLE_KEYWORDS if kw in line]
        if found_styles:
            style_info.append((line, found_styles))

    if style_info:
        for line, found in style_info:
            log(f"  [风格: {', '.join(found)}] → {line}")
    else:
        log(f"  (未找到明确的风格关键词)")
        for line in all_lines:
            if is_relevant_line(line):
                log(f"  {line[:200]}")

    # ---------- 尺寸分析 ----------
    log(f"\n▶ 2. 规格尺寸分析")
    dim_info = []
    for line in all_lines:
        found_dims = [kw for kw in DIM_KEYWORDS if kw in line]
        if found_dims:
            dim_info.append((line, found_dims))

    if dim_info:
        for line, found in dim_info:
            log(f"  [尺寸: {', '.join(found)}] → {line}")
    else:
        log(f"  (未找到明确尺寸关键词)")

    # ---------- 沙发类型 ----------
    log(f"\n▶ 3. 沙发类型/规格名称")
    sofa_info = []
    for line in all_lines:
        if any(kw in line for kw in SOFA_KEYWORDS):
            sofa_info.append(line)

    if sofa_info:
        seen = set()
        for line in sofa_info:
            if line not in seen:
                log(f"  → {line}")
                seen.add(line)
    else:
        log(f"  (未找到沙发类型信息)")

    # ---------- 带尺寸数字的文本 ----------
    log(f"\n▶ 4. 包含尺寸数字的文本")
    dim_pattern = re.compile(r'\d+[\s]*[×x\*][\s]*\d+')
    found_dim_lines = []
    for line in all_lines:
        if dim_pattern.search(line):
            found_dim_lines.append(line)
    if found_dim_lines:
        for line in found_dim_lines:
            log(f"  → {line}")
    else:
        log(f"  (未找到带尺寸数字的文本)")

    # ---------- 所有相关文本汇总 ----------
    log(f"\n▶ 5. 所有可能相关的文本行汇总")
    relevant_count = 0
    for line in all_lines:
        if is_relevant_line(line):
            relevant_count += 1
            if relevant_count <= 50:
                log(f"  {line[:300]}")
    if relevant_count > 50:
        log(f"  ...（共 {relevant_count} 行相关文本，仅显示前50行）")


def main():
    # 清空输出文件
    with open(OUTPUT_FILE, "w", encoding="utf-8") as f:
        f.write("")

    log("=" * 80)
    log("🛋️  顾家产品库 - 沙发PPT信息提取工具")
    log("=" * 80)

    for product_id, filename in FILES:
        analyze_product(product_id, filename)

    log(f"\n{'='*80}")
    log("✅ 所有文件分析完成！")
    log(f"{'='*80}")
    
    # 读取并打印结果
    with open(OUTPUT_FILE, "r", encoding="utf-8") as f:
        content = f.read()
    # Print first 8000 chars to terminal
    print(content[:8000])
    if len(content) > 8000:
        print(f"\n... [输出过长，完整内容已保存至 {OUTPUT_FILE}]")


if __name__ == "__main__":
    main()
